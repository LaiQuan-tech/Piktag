#!/usr/bin/env python3
"""Create PikTag presentation with screenshots and feature descriptions."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT_DIR = "/Users/aimand/.gemini/File/L PikTag/mobile/ppt-screenshots"
PPT_PATH = "/Users/aimand/.gemini/File/L PikTag/mobile/PikTag_功能介紹.pptx"

# Colors
GOLD = RGBColor(0xE8, 0xB8, 0x30)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
ACCENT = RGBColor(0xE8, 0x9D, 0x30)


def add_bg(slide, color=WHITE):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_phone_frame(slide, img_path, left, top, height):
    """Add a phone screenshot with rounded border effect."""
    # Add the screenshot image
    pic = slide.shapes.add_picture(img_path, left, top, height=height)
    return pic


def add_text_box(slide, text, left, top, width, height, font_size=14, bold=False, color=DARK, alignment=PP_ALIGN.LEFT):
    """Add a text box to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, items, left, top, width, height, font_size=12, color=DARK):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
        p.level = 0
    return txBox


def add_tag(slide, text, left, top, bg_color=GOLD, text_color=WHITE):
    """Add a tag/badge shape."""
    width = Inches(1.2)
    height = Inches(0.35)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = text_color
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============================================================
# Slide 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, WHITE)

# Title
add_text_box(slide, "# PikTag", Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             font_size=60, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)
# Subtitle
add_text_box(slide, "用標籤記住每個人", Inches(1), Inches(2.8), Inches(11), Inches(0.7),
             font_size=28, color=DARK, alignment=PP_ALIGN.CENTER)
# Description
add_text_box(slide, "行動社交 CRM — 功能總覽 & Demo", Inches(1), Inches(3.8), Inches(11), Inches(0.6),
             font_size=20, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# Version info
add_text_box(slide, "Phase 1 + Phase 2 + Phase 3  |  全功能已開發完成並部署", Inches(1), Inches(5.2), Inches(11), Inches(0.5),
             font_size=14, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# Tags
tags = ["#標籤管理", "#CRM提醒", "#社群連結", "#社交統計", "#附近搜尋"]
for i, tag in enumerate(tags):
    add_tag(slide, tag, Inches(2.5 + i * 1.7), Inches(6.0))

# ============================================================
# Slide 2: Feature Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, "功能架構總覽", Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             font_size=32, bold=True, color=DARK)

# Phase boxes
phases = [
    ("Phase 1 — 基礎功能", [
        "① 標籤式人脈管理",
        "② 個人檔案 & QR Code",
        "③ 搜尋與探索",
        "④ 即時聊天",
        "⑤ 便利貼筆記",
        "⑥ 社群連結 (Biolinks)",
    ], RGBColor(0xE3, 0xF2, 0xFD)),
    ("Phase 2 — CRM & 黏性", [
        "① Biolink 點擊追蹤通知",
        "② 生日/紀念日/合約到期提醒",
        "③ 歷史上的今天",
        "④ 標籤熱度通知",
    ], RGBColor(0xFF, 0xF8, 0xE1)),
    ("Phase 3 — 進階功能", [
        "① 附近熱門標籤推薦",
        "② 在這地點你認識誰",
        "③ 社交統計報表",
    ], RGBColor(0xFC, 0xE4, 0xEC)),
]

for i, (title, items, bg_c) in enumerate(phases):
    left = Inches(0.5 + i * 4.2)
    top = Inches(1.3)
    w = Inches(3.8)
    h = Inches(5.5)

    # Background box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_c
    shape.line.fill.background()

    # Phase title
    add_text_box(slide, title, left + Inches(0.3), top + Inches(0.2), w - Inches(0.6), Inches(0.5),
                 font_size=18, bold=True, color=DARK)

    # Items
    add_bullet_list(slide, items, left + Inches(0.3), top + Inches(0.9), w - Inches(0.6), h - Inches(1.2),
                    font_size=14, color=DARK)

# Status tag
add_tag(slide, "✅ 全部完成", Inches(10.5), Inches(0.4), bg_color=RGBColor(0x4C, 0xAF, 0x50))

# ============================================================
# Slide 3: Home Screen
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)

add_text_box(slide, "首頁 — 人脈總覽", Inches(0.5), Inches(0.3), Inches(5), Inches(0.7),
             font_size=32, bold=True, color=DARK)

# Phone screenshot
add_phone_frame(slide, f"{OUTPUT_DIR}/02_home.png", Inches(0.8), Inches(1.2), Inches(5.8))

# Feature description on the right
add_text_box(slide, "核心功能", Inches(6.5), Inches(1.2), Inches(6), Inches(0.5),
             font_size=22, bold=True, color=GOLD)

features = [
    "📅 動態日期標題 — 自動顯示今日日期",
    "📍 位置標籤 — 顯示目前所在區域",
    "🔔 通知鈴鐺 — 一鍵查看所有通知",
    "📌 定位圖示 — 快速查看附近人脈",
    "👤 人脈列表 — 顯示所有連結及標籤",
    "🏷️ 標籤系統 — 每人可加多個標籤分類",
    "⚙️ 篩選排序 — 支援多種排序方式",
]
add_bullet_list(slide, features, Inches(6.5), Inches(1.9), Inches(6), Inches(4.5), font_size=15)

add_text_box(slide, "💡 首頁也會顯示「今日提醒」和「歷史上的今天」卡片\n（當天有生日/紀念日的人脈時自動出現）",
             Inches(6.5), Inches(5.8), Inches(6), Inches(1), font_size=13, color=MEDIUM_GRAY)

# ============================================================
# Slide 4: Friend Detail
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)

add_text_box(slide, "好友詳情 — 完整人脈檔案", Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             font_size=32, bold=True, color=DARK)

# Three phone screenshots side by side
add_phone_frame(slide, f"{OUTPUT_DIR}/03_friend_detail_top.png", Inches(0.3), Inches(1.2), Inches(5.8))
add_phone_frame(slide, f"{OUTPUT_DIR}/04_friend_detail_crm.png", Inches(4.6), Inches(1.2), Inches(5.8))
add_phone_frame(slide, f"{OUTPUT_DIR}/05_friend_detail_biolinks.png", Inches(8.9), Inches(1.2), Inches(5.8))

# Labels under each screenshot
add_text_box(slide, "▲ 個人資料 & 標籤 & 相識紀錄", Inches(0.3), Inches(7.05), Inches(4), Inches(0.3),
             font_size=11, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, "▲ 便利貼 & CRM 重要提醒", Inches(4.6), Inches(7.05), Inches(4), Inches(0.3),
             font_size=11, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, "▲ 社群連結 (IG, LinkedIn)", Inches(8.9), Inches(7.05), Inches(4), Inches(0.3),
             font_size=11, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# Slide 5: CRM Features Detail
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, "Phase 2 — CRM & 黏性功能", Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             font_size=32, bold=True, color=DARK)

# CRM screenshot
add_phone_frame(slide, f"{OUTPUT_DIR}/04_friend_detail_crm.png", Inches(0.5), Inches(1.2), Inches(5.8))

# Feature details
add_text_box(slide, "CRM 重要提醒", Inches(6), Inches(1.2), Inches(6.5), Inches(0.5),
             font_size=22, bold=True, color=GOLD)

crm_features = [
    "🎂 生日提醒 — 記錄好友生日，到期自動通知",
    "💕 紀念日提醒 — 追蹤重要合作紀念日",
    "📋 合約到期提醒 — 商務人脈的合約管理",
    "",
    "📝 便利貼筆記",
    "• 為每位好友添加備忘便利貼",
    "• 支援釘選、編輯、刪除",
    "• 淡藍色卡片設計，一目了然",
    "",
    "🔗 社群連結追蹤",
    "• 自動記錄誰點擊了你的 Biolink",
    "• 觸發器自動建立通知",
    "• 追蹤 IG、LinkedIn 等平台點擊",
]
add_bullet_list(slide, crm_features, Inches(6), Inches(1.9), Inches(6.5), Inches(5), font_size=14)

# ============================================================
# Slide 6: Notifications
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)

add_text_box(slide, "通知中心 — 多維度通知", Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             font_size=32, bold=True, color=DARK)

# Notification screenshot
add_phone_frame(slide, f"{OUTPUT_DIR}/07_notifications.png", Inches(0.8), Inches(1.2), Inches(5.8))

# Feature description
add_text_box(slide, "四大分類標籤", Inches(6.5), Inches(1.2), Inches(6), Inches(0.5),
             font_size=22, bold=True, color=GOLD)

notif_features = [
    "📋 全部 — 一覽所有通知",
    "👥 追蹤 — 誰追蹤了你",
    "🏷️ 標籤 — 標籤熱度變化通知",
    "📊 CRM — Biolink 點擊 + 生日提醒",
    "",
    "五種通知類型：",
    "🔗 Biolink 點擊追蹤 — flowtest 點擊了你的 IG",
    "📈 標籤熱度上升 — #工程師 已有 88 人使用",
    "🎂 生日提醒 — 今天是小花的生日！",
    "👤 新追蹤者 — autotest2 開始追蹤你",
    "📅 歷史上的今天 — 1年前你認識了 Auto哥",
]
add_bullet_list(slide, notif_features, Inches(6.5), Inches(1.9), Inches(6), Inches(5), font_size=14)

# ============================================================
# Slide 7: Search Screen
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)

add_text_box(slide, "搜尋 — 多維度探索人脈", Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             font_size=32, bold=True, color=DARK)

# Search screenshot
add_phone_frame(slide, f"{OUTPUT_DIR}/06_search.png", Inches(0.8), Inches(1.2), Inches(5.8))

# Feature description
add_text_box(slide, "五大搜尋分類", Inches(6.5), Inches(1.2), Inches(6), Inches(0.5),
             font_size=22, bold=True, color=GOLD)

search_features = [
    "# 熱門標籤 — 瀏覽最多人使用的標籤",
    "📍 附近會員 — GPS 定位找附近的人",
    "✅ 認證會員 — 已驗證的專業人士",
    "🔥 附近熱標 (Phase 3 新增) — 你附近最火的標籤",
    "🕐 最近搜尋 — 快速回找",
    "",
    "熱門標籤排行：",
    "🥇 #工程師 — 88 位擁有",
    "🥈 #設計師 — 42 位擁有",
    "🥉 #創業家 — 33 位擁有",
    "④ #投資人 — 21 位擁有",
    "⑤ #台大校友 — 15 位擁有",
]
add_bullet_list(slide, search_features, Inches(6.5), Inches(1.9), Inches(6), Inches(5), font_size=14)

# ============================================================
# Slide 8: Social Stats
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)

add_text_box(slide, "Phase 3 — 社交統計報表", Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             font_size=32, bold=True, color=DARK)

# Two screenshots side by side
add_phone_frame(slide, f"{OUTPUT_DIR}/10_social_stats_top.png", Inches(0.5), Inches(1.2), Inches(5.8))
add_phone_frame(slide, f"{OUTPUT_DIR}/11_social_stats_bottom.png", Inches(4.5), Inches(1.2), Inches(5.8))

# Feature description on the right
add_text_box(slide, "數據一覽", Inches(9), Inches(1.2), Inches(4), Inches(0.5),
             font_size=22, bold=True, color=GOLD)

stats_features = [
    "📊 六大統計卡片",
    "• 總人脈數 & 月增長",
    "• 使用標籤數 & 平均值",
    "• 發送訊息數",
    "• 連結點擊數",
    "• 認證好友數",
    "• 便利貼數量",
    "",
    "📈 最常用標籤 Top 5",
    "• 視覺化長條圖排行",
    "",
    "📅 人脈時間軸",
    "• 最早 / 最新人脈日期",
    "",
    "📝 本週摘要",
    "• 自動統計本週動態",
    "",
    "⏱️ 時間篩選器",
    "• 本週 / 本月 / 全部",
]
add_bullet_list(slide, stats_features, Inches(9), Inches(1.9), Inches(4), Inches(5), font_size=13)

# ============================================================
# Slide 9: Settings
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)

add_text_box(slide, "設定 & 個人檔案", Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             font_size=32, bold=True, color=DARK)

# Two screenshots
add_phone_frame(slide, f"{OUTPUT_DIR}/08_profile.png", Inches(0.5), Inches(1.2), Inches(5.5))
add_phone_frame(slide, f"{OUTPUT_DIR}/09_settings.png", Inches(5), Inches(1.2), Inches(5.5))

# Features on the right
add_text_box(slide, "個人檔案", Inches(9.8), Inches(1.2), Inches(3.2), Inches(0.5),
             font_size=18, bold=True, color=GOLD)

profile_features = [
    "👤 頭像 & 用戶名",
    "📄 個人簡介",
    "🏷️ 標籤展示",
    "📤 分享個人檔案",
    "✏️ 編輯個人資訊",
    "📱 QR Code 名片",
]
add_bullet_list(slide, profile_features, Inches(9.8), Inches(1.8), Inches(3.2), Inches(2.5), font_size=12)

add_text_box(slide, "設定選項", Inches(9.8), Inches(4.2), Inches(3.2), Inches(0.5),
             font_size=18, bold=True, color=GOLD)

settings_features = [
    "📱 帳號資訊",
    "📞 通訊錄同步",
    "💌 邀請好友",
    "📍 在這地點你認識誰 ⭐",
    "📊 社交統計報表 ⭐",
    "🔒 隱私 / 通知設定",
    "🌐 語言 / 深色模式",
]
add_bullet_list(slide, settings_features, Inches(9.8), Inches(4.8), Inches(3.2), Inches(2.5), font_size=12)

# ============================================================
# Slide 10: Tech Stack & Testing
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, "技術架構 & 測試結果", Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             font_size=32, bold=True, color=DARK)

# Tech Stack Box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.5))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
shape.line.fill.background()

add_text_box(slide, "🛠️ 技術架構", Inches(0.8), Inches(1.5), Inches(5), Inches(0.5),
             font_size=20, bold=True, color=DARK)

tech = [
    "📱 前端：React Native (Expo SDK 54)",
    "☁️ 後端：Supabase (PostgreSQL + Auth)",
    "⚡ Edge Functions：Deno Runtime",
    "🔄 即時觸發器：PostgreSQL Triggers",
    "🔐 安全：RLS 政策 (40 條)",
    "🌐 部署：Vercel (Web Preview)",
    "📊 資料庫：15 張資料表",
    "",
    "🔧 關鍵技術亮點：",
    "• 觸發器自動建立通知 (biolink_click → notify)",
    "• 標籤快照追蹤熱度變化",
    "• GPS 定位附近人脈搜尋",
    "• Edge Function: daily-crm-check",
    "• Edge Function: suggest-tags",
]
add_bullet_list(slide, tech, Inches(0.8), Inches(2.1), Inches(5.2), Inches(4.5), font_size=13)

# Test Results Box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.3), Inches(5.8), Inches(5.5))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
shape.line.fill.background()

add_text_box(slide, "✅ 測試結果", Inches(7.1), Inches(1.5), Inches(5), Inches(0.5),
             font_size=20, bold=True, color=DARK)

tests = [
    "📋 資料庫驗證：29/29 全部通過",
    "• 8 張資料表存在性驗證",
    "• 3 個 CRM 欄位驗證",
    "• 7 項測試資料完整性",
    "• 2 個觸發器 + 2 個函式",
    "• 40 條 RLS 政策覆蓋",
    "• 端到端觸發器整合測試",
    "",
    "🖥️ UI 功能測試：22/22 全部通過",
    "• 首頁 5 項 / 搜尋 3 項",
    "• 好友詳情 5 項 / 通知 4 項",
    "• 設定 3 項 / 社交統計 2 項",
    "",
    "🔧 程式碼審查：14 個問題已修復",
    "• 4 個 Critical + 10 個 Major",
]
add_bullet_list(slide, tests, Inches(7.1), Inches(2.1), Inches(5.2), Inches(4.5), font_size=13)

# ============================================================
# Slide 11: Login Screen
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)

add_text_box(slide, "登入畫面", Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             font_size=32, bold=True, color=DARK)

add_phone_frame(slide, f"{OUTPUT_DIR}/01_login.png", Inches(4), Inches(1.0), Inches(5.5))

add_text_box(slide, "Supabase Auth 整合 — Email/Password 登入", Inches(2), Inches(6.8), Inches(9), Inches(0.5),
             font_size=16, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# Slide 12: Summary & Next Steps
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, "總結 & 下一步", Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             font_size=32, bold=True, color=DARK)

# Done box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(6), Inches(2.5))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
shape.line.fill.background()

add_text_box(slide, "✅ 已完成", Inches(0.8), Inches(1.5), Inches(5), Inches(0.4),
             font_size=20, bold=True, color=RGBColor(0x2E, 0x7D, 0x32))

done_items = [
    "Phase 1：6 大基礎功能",
    "Phase 2：4 大 CRM 黏性功能",
    "Phase 3：3 大進階功能",
    "全功能自動化測試 77 項通過",
    "14 個 Critical/Major 程式碼問題已修復",
]
add_bullet_list(slide, done_items, Inches(0.8), Inches(2.1), Inches(5.5), Inches(1.5), font_size=14)

# Next box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.3), Inches(6), Inches(2.5))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xE1)
shape.line.fill.background()

add_text_box(slide, "🔜 可優化項目 (12 Minor)", Inches(7.3), Inches(1.5), Inches(5), Inches(0.4),
             font_size=20, bold=True, color=ACCENT)

next_items = [
    "查詢效能優化 (Promise.all 並行化)",
    "FlatList 虛擬化改善",
    "Haversine 距離公式修正",
    "useMemo / useCallback 優化",
    "TypeScript 型別完整補齊",
]
add_bullet_list(slide, next_items, Inches(7.3), Inches(2.1), Inches(5.5), Inches(1.5), font_size=14)

# URLs
add_text_box(slide, "🌐 Live Demo", Inches(0.5), Inches(4.3), Inches(12), Inches(0.5),
             font_size=20, bold=True, color=DARK)
add_text_box(slide, "https://dist-gamma-pink.vercel.app", Inches(0.5), Inches(4.9), Inches(12), Inches(0.4),
             font_size=16, color=RGBColor(0x19, 0x76, 0xD2))

add_text_box(slide, "☁️ Supabase Project: kbwfdskulxnhjckdvghj", Inches(0.5), Inches(5.5), Inches(12), Inches(0.4),
             font_size=14, color=MEDIUM_GRAY)

# Thank you
add_text_box(slide, "Thank You 🙏", Inches(1), Inches(6.2), Inches(11), Inches(0.8),
             font_size=36, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)


# Save
prs.save(PPT_PATH)
print(f"✅ PPT saved to: {PPT_PATH}")
print(f"   Total slides: {len(prs.slides)}")
